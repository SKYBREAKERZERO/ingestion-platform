from __future__ import annotations

import math
from collections.abc import Iterator
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import (
    MSO_SHAPE_TYPE,
    PP_PLACEHOLDER,
)
from pptx.shapes.base import BaseShape
from pptx.shapes.shapetree import SlideShapes
from pptx.slide import Slide

from app.loader.base_loader import BaseLoader
from app.model.block import BlockType, DocumentBlock
from app.model.document import Document
from app.model.page import Page



def _enum_members(enum_type: Any, *names: str) -> set[Any]:
    """
    返回当前 python-pptx 版本实际存在的枚举成员。

    不同 python-pptx 版本的 PP_PLACEHOLDER 成员可能不同，
    因此禁止在类定义阶段直接访问可选成员。
    """

    members: set[Any] = set()

    for name in names:
        value = getattr(
            enum_type,
            name,
            None,
        )

        if value is not None:
            members.add(value)

    return members


class PPTXLoaderError(RuntimeError):
    """PPTX 文件加载异常。"""


class PPTXLoader(BaseLoader):
    """
    PPTX 原始内容加载器。

    默认映射：

        Presentation
            -> Document

        Slide
            -> Page

        Title / Subtitle
            -> HEADING Block

        TextBox / Body Placeholder
            -> PARAGRAPH、LIST 或 TEXTBOX Block

        Table Row
            -> TABLE Block

        Picture
            -> IMAGE Block

        Chart
            -> UNKNOWN Block，图表数据保存到 metadata

    负责：
        - 加载 PPTX
        - 读取 Slide 与 Shape
        - 提取标题、正文、列表、表格、图片和图表
        - 递归处理 Group Shape
        - 保留位置、尺寸、占位符、层级等元数据
        - 生成统一 Document、Page、DocumentBlock

    不负责：
        - OCR 图片
        - SmartArt 深层解析
        - Chapter / Section 识别
        - 内容过滤
        - Chunk
        - Token 统计
        - JSON / PostgreSQL 保存
    """

    SUPPORTED_EXTENSIONS = {
        ".pptx",
    }

    _TITLE_PLACEHOLDER_TYPES = _enum_members(
        PP_PLACEHOLDER,
        "TITLE",
        "CENTER_TITLE",
    )

    _SUBTITLE_PLACEHOLDER_TYPES = _enum_members(
        PP_PLACEHOLDER,
        "SUBTITLE",
    )

    _BODY_PLACEHOLDER_TYPES = _enum_members(
        PP_PLACEHOLDER,
        "BODY",
        "OBJECT",
        "TEXT",
    )

    def __init__(
        self,
        *,
        include_hidden_slides: bool = False,
        include_empty_slides: bool = True,
        include_images: bool = True,
        include_charts: bool = True,
        include_empty_image_blocks: bool = True,
        include_empty_chart_blocks: bool = True,
        extract_chart_data: bool = True,
        extract_text_per_paragraph: bool = True,
        extract_table_per_row: bool = True,
        recursive_group_shapes: bool = True,
        reading_order: str = "visual",
        cell_separator: str = " | ",
        paragraph_separator: str = "\n",
        maximum_slides: int | None = None,
        maximum_shapes_per_slide: int | None = None,
        maximum_table_rows: int | None = None,
        maximum_table_columns: int | None = None,
        print_summary: bool = False,
    ) -> None:

        if reading_order not in {
            "visual",
            "z_order",
        }:
            raise ValueError(
                "reading_order must be either "
                "'visual' or 'z_order'."
            )

        if not cell_separator:
            raise ValueError(
                "cell_separator cannot be empty."
            )

        if not paragraph_separator:
            raise ValueError(
                "paragraph_separator cannot be empty."
            )

        self._validate_optional_limit(
            "maximum_slides",
            maximum_slides,
        )

        self._validate_optional_limit(
            "maximum_shapes_per_slide",
            maximum_shapes_per_slide,
        )

        self._validate_optional_limit(
            "maximum_table_rows",
            maximum_table_rows,
        )

        self._validate_optional_limit(
            "maximum_table_columns",
            maximum_table_columns,
        )

        self.include_hidden_slides = (
            include_hidden_slides
        )

        self.include_empty_slides = (
            include_empty_slides
        )

        self.include_images = include_images
        self.include_charts = include_charts

        self.include_empty_image_blocks = (
            include_empty_image_blocks
        )

        self.include_empty_chart_blocks = (
            include_empty_chart_blocks
        )

        self.extract_chart_data = (
            extract_chart_data
        )

        self.extract_text_per_paragraph = (
            extract_text_per_paragraph
        )

        self.extract_table_per_row = (
            extract_table_per_row
        )

        self.recursive_group_shapes = (
            recursive_group_shapes
        )

        self.reading_order = reading_order
        self.cell_separator = cell_separator

        self.paragraph_separator = (
            paragraph_separator
        )

        self.maximum_slides = maximum_slides

        self.maximum_shapes_per_slide = (
            maximum_shapes_per_slide
        )

        self.maximum_table_rows = (
            maximum_table_rows
        )

        self.maximum_table_columns = (
            maximum_table_columns
        )

        self.print_summary = bool(
            print_summary
        )

    def load(
        self,
        file_path: str | Path,
    ) -> Document:

        path = self._validate_path(
            file_path
        )

        try:
            presentation = Presentation(
                str(path)
            )

            return self._build_document(
                path=path,
                presentation=presentation,
            )

        except PPTXLoaderError:
            raise

        except Exception as exc:
            raise PPTXLoaderError(
                f"Failed to load PPTX file "
                f"'{path.name}': {exc}"
            ) from exc

    def _build_document(
        self,
        *,
        path: Path,
        presentation: Any,
    ) -> Document:

        blocks: list[DocumentBlock] = []
        pages: list[Page] = []

        global_order = 0

        total_slide_count = len(
            presentation.slides
        )

        processed_slide_count = 0
        hidden_slide_count = 0
        skipped_slide_count = 0
        empty_slide_count = 0

        text_block_count = 0
        heading_block_count = 0
        list_block_count = 0
        table_block_count = 0
        image_block_count = 0
        chart_block_count = 0
        unknown_block_count = 0

        table_count = 0
        table_row_count = 0
        image_count = 0
        chart_count = 0
        group_shape_count = 0
        shape_count = 0

        slide_metadata: list[
            dict[str, Any]
        ] = []

        for slide_index, slide in enumerate(
            presentation.slides
        ):
            slide_number = slide_index + 1

            if (
                self.maximum_slides is not None
                and slide_number
                > self.maximum_slides
            ):
                break

            is_hidden = self._is_hidden_slide(
                slide
            )

            if is_hidden:
                hidden_slide_count += 1

                if not self.include_hidden_slides:
                    skipped_slide_count += 1

                    slide_metadata.append(
                        {
                            "slide_index": slide_index,
                            "slide_number": slide_number,
                            "hidden": True,
                            "status": "SKIPPED",
                            "reason": "hidden_slide",
                        }
                    )

                    continue

            processed_slide_count += 1

            slide_blocks: list[
                DocumentBlock
            ] = []

            slide_table_count = 0
            slide_image_count = 0
            slide_chart_count = 0
            slide_shape_count = 0

            shapes = list(
                self._iter_slide_shapes(
                    slide.shapes
                )
            )

            if self.reading_order == "visual":
                shapes = sorted(
                    shapes,
                    key=self._shape_visual_sort_key,
                )

            if (
                self.maximum_shapes_per_slide
                is not None
            ):
                shapes = shapes[
                    :self.maximum_shapes_per_slide
                ]

            for visual_index, shape_info in enumerate(
                shapes
            ):
                shape = shape_info["shape"]

                shape_index = int(
                    shape_info["shape_index"]
                )

                group_path = list(
                    shape_info["group_path"]
                )

                z_order = int(
                    shape_info["z_order"]
                )

                slide_shape_count += 1
                shape_count += 1

                if group_path:
                    group_shape_count += 1

                created_blocks = (
                    self._extract_shape_blocks(
                        shape=shape,
                        slide_index=slide_index,
                        slide_number=slide_number,
                        shape_index=shape_index,
                        visual_index=visual_index,
                        z_order=z_order,
                        group_path=group_path,
                        start_order=global_order,
                    )
                )

                for block in created_blocks:
                    slide_blocks.append(block)
                    blocks.append(block)

                    global_order += 1

                    if block.block_type == BlockType.HEADING:
                        heading_block_count += 1

                    elif block.block_type in {
                        BlockType.PARAGRAPH,
                        BlockType.TEXTBOX,
                    }:
                        text_block_count += 1

                    elif block.block_type == BlockType.LIST:
                        list_block_count += 1

                    elif block.block_type == BlockType.TABLE:
                        table_block_count += 1
                        table_row_count += 1

                    elif block.block_type == BlockType.IMAGE:
                        image_block_count += 1

                    else:
                        unknown_block_count += 1

                    if block.metadata.get(
                        "content_kind"
                    ) == "chart":
                        chart_block_count += 1

                if self._shape_has_table(shape):
                    table_count += 1
                    slide_table_count += 1

                if self._is_picture_shape(shape):
                    image_count += 1
                    slide_image_count += 1

                if self._shape_has_chart(shape):
                    chart_count += 1
                    slide_chart_count += 1

            slide_text = self._build_slide_text(
                slide_blocks
            )

            if not slide_text:
                empty_slide_count += 1

                if not self.include_empty_slides:
                    for block in slide_blocks:
                        if block in blocks:
                            blocks.remove(block)

                    global_order = len(blocks)
                    processed_slide_count -= 1
                    skipped_slide_count += 1

                    slide_metadata.append(
                        {
                            "slide_index": slide_index,
                            "slide_number": slide_number,
                            "hidden": is_hidden,
                            "status": "SKIPPED",
                            "reason": "empty_slide",
                        }
                    )

                    continue

            logical_page_number = (
                len(pages) + 1
            )

            for block in slide_blocks:
                block.page_number = (
                    logical_page_number
                )

                block.metadata[
                    "logical_page_number"
                ] = logical_page_number

            pages.append(
                Page(
                    page_number=logical_page_number,
                    text=slide_text,
                )
            )

            slide_title = self._resolve_slide_title(
                slide_blocks
            )

            slide_metadata.append(
                {
                    "slide_index": slide_index,
                    "slide_number": slide_number,
                    "logical_page_number": (
                        logical_page_number
                    ),
                    "title": slide_title,
                    "hidden": is_hidden,
                    "status": (
                        "EMPTY"
                        if not slide_text
                        else "SUCCESS"
                    ),
                    "shape_count": slide_shape_count,
                    "block_count": len(
                        slide_blocks
                    ),
                    "table_count": slide_table_count,
                    "image_count": slide_image_count,
                    "chart_count": slide_chart_count,
                    "character_count": len(
                        slide_text
                    ),
                }
            )

        character_count = sum(
            len(page.text)
            for page in pages
        )

        metadata: dict[str, Any] = {
            "source_format": "pptx",
            "loader": "PPTXLoader",
            "loader_status": "SUCCESS",
            "slide_width_emu": int(
                presentation.slide_width
            ),
            "slide_height_emu": int(
                presentation.slide_height
            ),
            "slide_width_inches": (
                self._emu_to_inches(
                    presentation.slide_width
                )
            ),
            "slide_height_inches": (
                self._emu_to_inches(
                    presentation.slide_height
                )
            ),
            "slide_count": total_slide_count,
            "processed_slide_count": (
                processed_slide_count
            ),
            "skipped_slide_count": (
                skipped_slide_count
            ),
            "hidden_slide_count": (
                hidden_slide_count
            ),
            "empty_slide_count": (
                empty_slide_count
            ),
            "shape_count": shape_count,
            "group_shape_count": (
                group_shape_count
            ),
            "block_count": len(blocks),
            "heading_block_count": (
                heading_block_count
            ),
            "text_block_count": (
                text_block_count
            ),
            "list_block_count": (
                list_block_count
            ),
            "table_count": table_count,
            "table_row_count": (
                table_row_count
            ),
            "table_block_count": (
                table_block_count
            ),
            "image_count": image_count,
            "image_block_count": (
                image_block_count
            ),
            "chart_count": chart_count,
            "chart_block_count": (
                chart_block_count
            ),
            "unknown_block_count": (
                unknown_block_count
            ),
            "character_count": character_count,
            "reading_order": self.reading_order,
            "slides": slide_metadata,
        }

        if self.print_summary:
            self._print_summary(
                path=path,
                metadata=metadata,
            )

        return Document(
            file_name=path.name,
            file_type="pptx",
            pages=pages,
            blocks=blocks,
            chapters=[],
            sections=[],
            contents=[],
            metadata=metadata,
        )

    # ==================================================
    # Shape extraction
    # ==================================================

    def _extract_shape_blocks(
        self,
        *,
        shape: BaseShape,
        slide_index: int,
        slide_number: int,
        shape_index: int,
        visual_index: int,
        z_order: int,
        group_path: list[int],
        start_order: int,
    ) -> list[DocumentBlock]:

        common_metadata = self._build_shape_metadata(
            shape=shape,
            slide_index=slide_index,
            slide_number=slide_number,
            shape_index=shape_index,
            visual_index=visual_index,
            z_order=z_order,
            group_path=group_path,
        )

        if self._shape_has_table(shape):
            return self._extract_table_blocks(
                shape=shape,
                metadata=common_metadata,
                start_order=start_order,
            )

        if self._shape_has_chart(shape):
            if not self.include_charts:
                return []

            return self._extract_chart_blocks(
                shape=shape,
                metadata=common_metadata,
                order=start_order,
            )

        if self._is_picture_shape(shape):
            if not self.include_images:
                return []

            return self._extract_image_blocks(
                shape=shape,
                metadata=common_metadata,
                order=start_order,
            )

        if self._shape_has_text(shape):
            return self._extract_text_blocks(
                shape=shape,
                metadata=common_metadata,
                start_order=start_order,
            )

        return []

    def _extract_text_blocks(
        self,
        *,
        shape: BaseShape,
        metadata: dict[str, Any],
        start_order: int,
    ) -> list[DocumentBlock]:

        text_frame = getattr(
            shape,
            "text_frame",
            None,
        )

        if text_frame is None:
            return []

        placeholder_type = (
            self._resolve_placeholder_type(
                shape
            )
        )

        is_title = (
            placeholder_type
            in self._TITLE_PLACEHOLDER_TYPES
        )

        is_subtitle = (
            placeholder_type
            in self._SUBTITLE_PLACEHOLDER_TYPES
        )

        blocks: list[
            DocumentBlock
        ] = []

        paragraphs = list(
            text_frame.paragraphs
        )

        if not self.extract_text_per_paragraph:
            normalized_paragraphs = [
                self._normalize_text(
                    paragraph.text
                )
                for paragraph in paragraphs
            ]

            combined_text = (
                self.paragraph_separator.join(
                    text
                    for text
                    in normalized_paragraphs
                    if text
                )
            ).strip()

            if not combined_text:
                return []

            block_type, level = (
                self._resolve_text_block_type(
                    is_title=is_title,
                    is_subtitle=is_subtitle,
                    paragraph_level=0,
                    shape=shape,
                )
            )

            blocks.append(
                DocumentBlock(
                    id=self._build_block_id(
                        metadata=metadata,
                        suffix="text",
                    ),
                    block_type=block_type,
                    text=combined_text,
                    level=level,
                    style_name=metadata.get(
                        "placeholder_type"
                    ),
                    order=start_order,
                    page_number=metadata[
                        "slide_number"
                    ],
                    source="pptx",
                    metadata={
                        **metadata,
                        "content_kind": "text",
                        "paragraph_count": len(
                            paragraphs
                        ),
                    },
                )
            )

            return blocks

        local_order = start_order

        for paragraph_index, paragraph in enumerate(
            paragraphs
        ):
            text = self._normalize_text(
                paragraph.text
            )

            if not text:
                continue

            paragraph_level = int(
                getattr(
                    paragraph,
                    "level",
                    0,
                )
                or 0
            )

            block_type, level = (
                self._resolve_text_block_type(
                    is_title=is_title,
                    is_subtitle=is_subtitle,
                    paragraph_level=(
                        paragraph_level
                    ),
                    shape=shape,
                )
            )

            blocks.append(
                DocumentBlock(
                    id=self._build_block_id(
                        metadata=metadata,
                        suffix=(
                            f"paragraph-"
                            f"{paragraph_index}"
                        ),
                    ),
                    block_type=block_type,
                    text=text,
                    level=level,
                    style_name=metadata.get(
                        "placeholder_type"
                    ),
                    order=local_order,
                    page_number=metadata[
                        "slide_number"
                    ],
                    source="pptx",
                    metadata={
                        **metadata,
                        "content_kind": "text",
                        "paragraph_index": (
                            paragraph_index
                        ),
                        "paragraph_level": (
                            paragraph_level
                        ),
                        "run_count": len(
                            paragraph.runs
                        ),
                    },
                )
            )

            local_order += 1

        return blocks

    def _extract_table_blocks(
        self,
        *,
        shape: BaseShape,
        metadata: dict[str, Any],
        start_order: int,
    ) -> list[DocumentBlock]:

        table = shape.table

        blocks: list[
            DocumentBlock
        ] = []

        row_limit = len(
            table.rows
        )

        if self.maximum_table_rows is not None:
            row_limit = min(
                row_limit,
                self.maximum_table_rows,
            )

        local_order = start_order

        for row_index in range(row_limit):
            row = table.rows[
                row_index
            ]

            column_limit = len(
                row.cells
            )

            if (
                self.maximum_table_columns
                is not None
            ):
                column_limit = min(
                    column_limit,
                    self.maximum_table_columns,
                )

            cells: list[str] = []

            for column_index in range(
                column_limit
            ):
                cell = row.cells[
                    column_index
                ]

                value = self._normalize_text(
                    cell.text
                )

                cells.append(value)

            # Loader 层保留所有 Cell 的列位置。
            #
            # 例如：
            #
            #     ["", "A", "", "C"]
            #
            # 不能在这里裁剪成：
            #
            #     ["A", "", "C"]
            #
            # 否则后续 TableFilter / Parser 已无法恢复原始列结构。
            if not any(cells):
                continue

            row_text = self.cell_separator.join(
                cells
            )

            blocks.append(
                DocumentBlock(
                    id=self._build_block_id(
                        metadata=metadata,
                        suffix=(
                            f"table-row-{row_index}"
                        ),
                    ),
                    block_type=BlockType.TABLE,
                    text=row_text,
                    order=local_order,
                    page_number=metadata[
                        "slide_number"
                    ],
                    table_index=metadata[
                        "shape_index"
                    ],
                    row_index=row_index,
                    cells=cells,
                    source="pptx",
                    metadata={
                        **metadata,
                        "content_kind": "table",
                        "table_row_index": (
                            row_index
                        ),
                        "table_row_count": len(
                            table.rows
                        ),
                        "table_column_count": (
                            column_limit
                        ),
                        "non_empty_cell_count": (
                            sum(
                                1
                                for cell
                                in cells
                                if cell
                            )
                        ),
                        "column_position_preserved": True,
                        "is_header_candidate": (
                            row_index == 0
                        ),
                    },
                )
            )

            local_order += 1

        return blocks

    def _extract_image_blocks(
        self,
        *,
        shape: BaseShape,
        metadata: dict[str, Any],
        order: int,
    ) -> list[DocumentBlock]:

        image_metadata = {
            **metadata,
            "content_kind": "image",
            "has_image": True,
        }

        image = getattr(
            shape,
            "image",
            None,
        )

        if image is not None:
            image_metadata.update(
                {
                    "image_filename": getattr(
                        image,
                        "filename",
                        None,
                    ),
                    "image_extension": getattr(
                        image,
                        "ext",
                        None,
                    ),
                    "image_content_type": getattr(
                        image,
                        "content_type",
                        None,
                    ),
                    "image_size_bytes": len(
                        image.blob
                    ),
                }
            )

        alt_text = self._extract_alt_text(
            shape
        )

        image_name = self._normalize_text(
            getattr(
                shape,
                "name",
                "",
            )
        )

        image_text = (
            alt_text
            or image_name
        )

        if (
            not image_text
            and not self.include_empty_image_blocks
        ):
            return []

        return [
            DocumentBlock(
                id=self._build_block_id(
                    metadata=metadata,
                    suffix="image",
                ),
                block_type=BlockType.IMAGE,
                text=image_text,
                order=order,
                page_number=metadata[
                    "slide_number"
                ],
                source="pptx",
                metadata=image_metadata,
            )
        ]

    def _extract_chart_blocks(
        self,
        *,
        shape: BaseShape,
        metadata: dict[str, Any],
        order: int,
    ) -> list[DocumentBlock]:

        chart_metadata: dict[str, Any] = {
            **metadata,
            "content_kind": "chart",
            "has_chart": True,
        }

        chart = shape.chart

        chart_title = self._extract_chart_title(
            chart
        )

        if self.extract_chart_data:
            chart_metadata.update(
                self._extract_chart_data(
                    chart
                )
            )

        text_parts: list[str] = []

        if chart_title:
            text_parts.append(
                chart_title
            )

        series_names = chart_metadata.get(
            "chart_series_names",
            [],
        )

        if series_names:
            text_parts.append(
                "Series: "
                + ", ".join(
                    str(name)
                    for name in series_names
                )
            )

        categories = chart_metadata.get(
            "chart_categories",
            [],
        )

        if categories:
            text_parts.append(
                "Categories: "
                + ", ".join(
                    str(category)
                    for category in categories
                )
            )

        chart_text = self.paragraph_separator.join(
            text_parts
        ).strip()

        if (
            not chart_text
            and not self.include_empty_chart_blocks
        ):
            return []

        # 当前 BlockType 没有 CHART，因此使用 UNKNOWN，
        # 并由 metadata["content_kind"] 标识为 chart。
        return [
            DocumentBlock(
                id=self._build_block_id(
                    metadata=metadata,
                    suffix="chart",
                ),
                block_type=BlockType.UNKNOWN,
                text=chart_text,
                order=order,
                page_number=metadata[
                    "slide_number"
                ],
                source="pptx",
                metadata=chart_metadata,
            )
        ]

    # ==================================================
    # Group shape
    # ==================================================

    def _iter_slide_shapes(
        self,
        shapes: SlideShapes,
    ) -> Iterator[dict[str, Any]]:

        for z_order, shape in enumerate(
            shapes
        ):
            yield from self._walk_shape(
                shape=shape,
                shape_index=z_order,
                z_order=z_order,
                group_path=[],
            )

    def _walk_shape(
        self,
        *,
        shape: BaseShape,
        shape_index: int,
        z_order: int,
        group_path: list[int],
    ) -> Iterator[dict[str, Any]]:

        if (
            self.recursive_group_shapes
            and self._is_group_shape(shape)
        ):
            group_shape = shape

            for child_index, child in enumerate(
                group_shape.shapes
            ):
                yield from self._walk_shape(
                    shape=child,
                    shape_index=child_index,
                    z_order=child_index,
                    group_path=[
                        *group_path,
                        shape_index,
                    ],
                )

            return

        yield {
            "shape": shape,
            "shape_index": shape_index,
            "z_order": z_order,
            "group_path": group_path,
        }

    # ==================================================
    # Chart
    # ==================================================

    @classmethod
    def _extract_chart_data(
        cls,
        chart,
    ) -> dict[str, Any]:

        result: dict[str, Any] = {
            "chart_type": str(
                getattr(
                    chart,
                    "chart_type",
                    "",
                )
            ),
            "chart_series_names": [],
            "chart_categories": [],
            "chart_series": [],
        }

        try:
            plots = list(
                chart.plots
            )

            if plots:
                categories = getattr(
                    plots[0],
                    "categories",
                    None,
                )

                if categories is not None:
                    result[
                        "chart_categories"
                    ] = [
                        cls._normalize_chart_value(
                            category.label
                            if hasattr(
                                category,
                                "label",
                            )
                            else category
                        )
                        for category in categories
                    ]

        except Exception:
            pass

        try:
            series_data: list[
                dict[str, Any]
            ] = []

            series_names: list[str] = []

            for series_index, series in enumerate(
                chart.series
            ):
                name = cls._normalize_text(
                    getattr(
                        series,
                        "name",
                        "",
                    )
                )

                if name:
                    series_names.append(name)

                values: list[Any] = []

                try:
                    values = [
                        cls._normalize_chart_value(
                            value
                        )
                        for value in series.values
                    ]

                except Exception:
                    values = []

                series_data.append(
                    {
                        "series_index": series_index,
                        "name": name or None,
                        "values": values,
                    }
                )

            result[
                "chart_series_names"
            ] = series_names

            result[
                "chart_series"
            ] = series_data

        except Exception:
            pass

        return result

    @classmethod
    def _extract_chart_title(
        cls,
        chart,
    ) -> str:

        try:
            if not chart.has_title:
                return ""

            title = chart.chart_title

            if not title.has_text_frame:
                return ""

            return cls._normalize_text(
                title.text_frame.text
            )

        except Exception:
            return ""

    @staticmethod
    def _normalize_chart_value(
        value: Any,
    ) -> Any:

        if value is None:
            return None

        if isinstance(value, bool):
            return value

        if isinstance(value, int):
            return value

        if isinstance(value, float):
            if math.isnan(value):
                return None

            if math.isinf(value):
                return None

            return value

        label = getattr(
            value,
            "label",
            None,
        )

        if label is not None:
            return str(label)

        return str(value)

    # ==================================================
    # Metadata
    # ==================================================

    @classmethod
    def _build_shape_metadata(
        cls,
        *,
        shape: BaseShape,
        slide_index: int,
        slide_number: int,
        shape_index: int,
        visual_index: int,
        z_order: int,
        group_path: list[int],
    ) -> dict[str, Any]:

        placeholder_type = (
            cls._resolve_placeholder_type(
                shape
            )
        )

        shape_type = getattr(
            shape,
            "shape_type",
            None,
        )

        return {
            "source": "pptx",
            "slide_index": slide_index,
            "slide_number": slide_number,
            "shape_index": shape_index,
            "visual_index": visual_index,
            "z_order": z_order,
            "group_path": group_path,
            "shape_id": getattr(
                shape,
                "shape_id",
                None,
            ),
            "shape_name": getattr(
                shape,
                "name",
                None,
            ),
            "shape_type": (
                str(shape_type)
                if shape_type is not None
                else None
            ),
            "placeholder_type": (
                cls._enum_name(
                    placeholder_type
                )
                if placeholder_type is not None
                else None
            ),
            "is_placeholder": bool(
                getattr(
                    shape,
                    "is_placeholder",
                    False,
                )
            ),
            "left_emu": cls._safe_int(
                getattr(
                    shape,
                    "left",
                    None,
                )
            ),
            "top_emu": cls._safe_int(
                getattr(
                    shape,
                    "top",
                    None,
                )
            ),
            "width_emu": cls._safe_int(
                getattr(
                    shape,
                    "width",
                    None,
                )
            ),
            "height_emu": cls._safe_int(
                getattr(
                    shape,
                    "height",
                    None,
                )
            ),
            "left_inches": cls._emu_to_inches(
                getattr(
                    shape,
                    "left",
                    None,
                )
            ),
            "top_inches": cls._emu_to_inches(
                getattr(
                    shape,
                    "top",
                    None,
                )
            ),
            "width_inches": cls._emu_to_inches(
                getattr(
                    shape,
                    "width",
                    None,
                )
            ),
            "height_inches": cls._emu_to_inches(
                getattr(
                    shape,
                    "height",
                    None,
                )
            ),
            "rotation": getattr(
                shape,
                "rotation",
                None,
            ),
        }

    @staticmethod
    def _build_block_id(
        *,
        metadata: dict[str, Any],
        suffix: str,
    ) -> str:

        slide_index = int(
            metadata["slide_index"]
        )

        shape_index = int(
            metadata["shape_index"]
        )

        group_path = metadata.get(
            "group_path",
            [],
        )

        group_part = ""

        if group_path:
            group_part = (
                "-group-"
                + "-".join(
                    str(value)
                    for value in group_path
                )
            )

        return (
            f"pptx-slide-{slide_index:04d}"
            f"{group_part}"
            f"-shape-{shape_index:04d}"
            f"-{suffix}"
        )

    # ==================================================
    # Shape type checks
    # ==================================================

    @staticmethod
    def _shape_has_text(
        shape: BaseShape,
    ) -> bool:

        try:
            return bool(
                getattr(
                    shape,
                    "has_text_frame",
                    False,
                )
            )

        except Exception:
            return False

    @staticmethod
    def _shape_has_table(
        shape: BaseShape,
    ) -> bool:

        try:
            return bool(
                getattr(
                    shape,
                    "has_table",
                    False,
                )
            )

        except Exception:
            return False

    @staticmethod
    def _shape_has_chart(
        shape: BaseShape,
    ) -> bool:

        try:
            return bool(
                getattr(
                    shape,
                    "has_chart",
                    False,
                )
            )

        except Exception:
            return False

    @staticmethod
    def _is_picture_shape(
        shape: BaseShape,
    ) -> bool:

        shape_type = getattr(
            shape,
            "shape_type",
            None,
        )

        picture_types = {
            MSO_SHAPE_TYPE.PICTURE,
        }

        linked_picture = getattr(
            MSO_SHAPE_TYPE,
            "LINKED_PICTURE",
            None,
        )

        if linked_picture is not None:
            picture_types.add(
                linked_picture
            )

        return shape_type in picture_types

    @staticmethod
    def _is_group_shape(
        shape: BaseShape,
    ) -> bool:

        return (
            getattr(
                shape,
                "shape_type",
                None,
            )
            == MSO_SHAPE_TYPE.GROUP
        )

    # ==================================================
    # Text classification
    # ==================================================

    @classmethod
    def _resolve_text_block_type(
        cls,
        *,
        is_title: bool,
        is_subtitle: bool,
        paragraph_level: int,
        shape: BaseShape,
    ) -> tuple[BlockType, int | None]:

        if is_title:
            return BlockType.HEADING, 1

        if is_subtitle:
            return BlockType.HEADING, 2

        if paragraph_level > 0:
            return (
                BlockType.LIST,
                paragraph_level + 1,
            )

        placeholder_type = (
            cls._resolve_placeholder_type(
                shape
            )
        )

        if (
            placeholder_type
            in cls._BODY_PLACEHOLDER_TYPES
        ):
            return BlockType.PARAGRAPH, None

        return BlockType.TEXTBOX, None

    @staticmethod
    def _resolve_placeholder_type(
        shape: BaseShape,
    ):

        try:
            if not shape.is_placeholder:
                return None

            return shape.placeholder_format.type

        except Exception:
            return None

    # ==================================================
    # Reading order
    # ==================================================

    @staticmethod
    def _shape_visual_sort_key(
        shape_info: dict[str, Any],
    ) -> tuple[int, int, int]:

        shape = shape_info["shape"]

        top = int(
            getattr(
                shape,
                "top",
                0,
            )
            or 0
        )

        left = int(
            getattr(
                shape,
                "left",
                0,
            )
            or 0
        )

        z_order = int(
            shape_info["z_order"]
        )

        return (
            top,
            left,
            z_order,
        )

    # ==================================================
    # Slide helpers
    # ==================================================

    @staticmethod
    def _is_hidden_slide(
        slide: Slide,
    ) -> bool:
        """
        检查 PowerPoint 的 show="0" 标志。

        python-pptx 当前没有统一高层 hidden 属性，
        因此仅在此处读取底层 XML 属性。
        """

        try:
            value = slide._element.get(
                "show"
            )

            return value in {
                "0",
                "false",
                "False",
            }

        except Exception:
            return False

    @staticmethod
    def _resolve_slide_title(
        blocks: list[DocumentBlock],
    ) -> str | None:

        for block in blocks:
            if (
                block.block_type
                == BlockType.HEADING
                and block.level == 1
                and block.text
            ):
                return block.text

        for block in blocks:
            if block.text:
                return block.text[
                    :200
                ]

        return None

    @staticmethod
    def _build_slide_text(
        blocks: list[DocumentBlock],
    ) -> str:

        return "\n".join(
            block.text.strip()
            for block in blocks
            if block.text.strip()
        ).strip()

    # ==================================================
    # Image helpers
    # ==================================================

    @classmethod
    def _extract_alt_text(
        cls,
        shape: BaseShape,
    ) -> str:

        try:
            element = shape._element

            c_nv_pr = element.xpath(
                ".//*[local-name()='cNvPr']"
            )

            if not c_nv_pr:
                return ""

            description = c_nv_pr[0].get(
                "descr"
            )

            title = c_nv_pr[0].get(
                "title"
            )

            return cls._normalize_text(
                description
                or title
                or ""
            )

        except Exception:
            return ""

    # ==================================================
    # General helpers
    # ==================================================

    @staticmethod
    def _trim_empty_boundaries(
        cells: list[str],
    ) -> list[str]:

        if not cells:
            return []

        start = 0
        end = len(cells)

        while (
            start < end
            and not cells[start]
        ):
            start += 1

        while (
            end > start
            and not cells[end - 1]
        ):
            end -= 1

        return cells[
            start:end
        ]

    @staticmethod
    def _normalize_text(
        value: Any,
    ) -> str:

        if value is None:
            return ""

        normalized = str(
            value
        )

        normalized = normalized.replace(
            "\u3000",
            " ",
        )

        normalized = normalized.replace(
            "\xa0",
            " ",
        )

        for character in (
            "\u200b",
            "\u200c",
            "\u200d",
            "\u2060",
            "\ufeff",
        ):
            normalized = normalized.replace(
                character,
                "",
            )

        normalized = normalized.replace(
            "\r\n",
            "\n",
        )

        normalized = normalized.replace(
            "\r",
            "\n",
        )

        lines = [
            " ".join(
                line.split()
            )
            for line in normalized.splitlines()
            if line.strip()
        ]

        return "\n".join(
            lines
        ).strip()

    @staticmethod
    def _enum_name(
        value: Any,
    ) -> str:

        name = getattr(
            value,
            "name",
            None,
        )

        if name:
            return str(name)

        return str(value)

    @staticmethod
    def _safe_int(
        value: Any,
    ) -> int | None:

        if value is None:
            return None

        try:
            return int(value)

        except (
            TypeError,
            ValueError,
        ):
            return None

    @staticmethod
    def _emu_to_inches(
        value: Any,
    ) -> float | None:

        if value is None:
            return None

        try:
            return round(
                int(value) / 914400,
                4,
            )

        except (
            TypeError,
            ValueError,
        ):
            return None

    @staticmethod
    def _validate_optional_limit(
        name: str,
        value: int | None,
    ) -> None:

        if value is not None and value <= 0:
            raise ValueError(
                f"{name} must be greater than 0."
            )

    @classmethod
    def _validate_path(
        cls,
        file_path: str | Path,
    ) -> Path:

        if file_path is None:
            raise ValueError(
                "file_path cannot be None."
            )

        if not str(
            file_path
        ).strip():
            raise ValueError(
                "file_path cannot be empty."
            )

        path = Path(
            file_path
        ).expanduser()

        if not path.exists():
            raise FileNotFoundError(
                f"PPTX file not found: {path}"
            )

        if not path.is_file():
            raise IsADirectoryError(
                f"Input path is not a file: {path}"
            )

        if path.name.startswith("~$"):
            raise ValueError(
                "Temporary PowerPoint file is not supported: "
                f"{path.name}"
            )

        if (
            path.suffix.lower()
            not in cls.SUPPORTED_EXTENSIONS
        ):
            supported = ", ".join(
                sorted(
                    cls.SUPPORTED_EXTENSIONS
                )
            )

            raise ValueError(
                f"PPTXLoader only accepts: "
                f"{supported}. "
                f"Received: "
                f"{path.suffix or '<no extension>'}"
            )

        return path

    @staticmethod
    def _print_summary(
        *,
        path: Path,
        metadata: dict[str, Any],
    ) -> None:

        print()
        print("===== PPTX Loader =====")
        print(
            f"File             : {path.name}"
        )
        print(
            "Slides           :",
            metadata["slide_count"],
        )
        print(
            "Processed slides :",
            metadata[
                "processed_slide_count"
            ],
        )
        print(
            "Skipped slides   :",
            metadata[
                "skipped_slide_count"
            ],
        )
        print(
            "Hidden slides    :",
            metadata[
                "hidden_slide_count"
            ],
        )
        print(
            "Empty slides     :",
            metadata[
                "empty_slide_count"
            ],
        )
        print(
            "Shapes           :",
            metadata["shape_count"],
        )
        print(
            "Blocks           :",
            metadata["block_count"],
        )
        print(
            "Headings         :",
            metadata[
                "heading_block_count"
            ],
        )
        print(
            "Text blocks      :",
            metadata[
                "text_block_count"
            ],
        )
        print(
            "List blocks      :",
            metadata[
                "list_block_count"
            ],
        )
        print(
            "Tables           :",
            metadata["table_count"],
        )
        print(
            "Table rows       :",
            metadata[
                "table_row_count"
            ],
        )
        print(
            "Images           :",
            metadata["image_count"],
        )
        print(
            "Charts           :",
            metadata["chart_count"],
        )
        print(
            "Characters       :",
            metadata[
                "character_count"
            ],
        )
        print("=======================")
        print()